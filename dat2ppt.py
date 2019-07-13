# -*- coding:utf-8 -*-
# Author:XuQiang
# Data:2019/7/11 10:00
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import time
import re
import os


def read_input():
    try:
        with open('Input.dat', 'r') as txt_f:
            for line in txt_f:
                path = re.split('[=,]', line)
                while '' in path:
                    path.remove('')
                for i in range(len(path)):
                    path[i] = path[i].strip()
                    path[i] = path[i].strip('\'')
                    path[i] = path[i].strip('\"')  # 删除可能存在的空格、引号
                    if '\\' in path[i]:
                        path[i] = path[i].replace('\\', '/')
                if 'InstallPath' in path:
                    install_path = path[1]  # 安装目录
                if 'WorkPath' in path:
                    work_path = path[1]  # 工作目录
                if 'ImageName' in path:
                    image_name = path  # 图片名字
                if 'ReportName' in path:
                    report_name = path  # ppt报告
                if 'TxtName' in path:
                    txt_name = path  # ppt内容
        return install_path, work_path, image_name, report_name, txt_name
    except NameError:
        return False


def read_txt(txt_filepath):
    # 读取txt文件
    index = [0] * 5
    word = ['', '', '', '', '', '']
    with open(txt_filepath, 'r') as f:
        line_index = 0
        for line in f:
            if not line.isspace():
                if '优化信息' in line:
                    index[0] = line_index
                elif '轮胎参数' in line:
                    index[1] = line_index
                elif '节距参数' in line:
                    index[2] = line_index
                elif '初始排列' in line:
                    index[3] = line_index
                elif '优化排列' in line:
                    index[4] = line_index
                elif '初始能量幅值' in line:
                    init_line = line.split(' ')
                    while '' in init_line:
                        init_line.remove('')
                    word[3] = init_line[1].strip()
                elif '最优能量幅值' in line:
                    init_line = line.split(' ')
                    while '' in init_line:
                        init_line.remove('')
                    word[4] = init_line[1].strip()
            line_index += 1
    with open(txt_filepath, 'r') as f1:
        for lin in f1.readlines()[index[0] + 1:index[1] - 1]:
            if '保存路径' not in lin:
                word[0] += lin.split(':')[-1].strip(' ')
    with open(txt_filepath, 'r') as f2:
        for lin in f2.readlines()[index[1]:index[2]-1]:
            if len(lin.strip()) > 1:
                word[1] += lin.strip(' ')
    with open(txt_filepath, 'r') as f3:
        for lin in f3.readlines()[index[2]:index[3]-1]:
            if len(lin.strip()) > 1:
                word[1] += lin.strip(' ')
    with open(txt_filepath, 'r') as f4:
        for lin in f4.readlines()[index[3]:index[4]-1]:
            if len(lin.strip()) > 1:
                word[2] += lin.strip(' ')
    with open(txt_filepath, 'r') as f5:
        for lin in f5.readlines()[index[4]:]:
            if len(lin.strip()) > 1:
                word[2] += lin.strip(' ')
    word[5] += txt_filepath.split('/')[-1][:-4] + '\n'
    return word


def create_ppt(template, img_path, report_path, word):
    # 导入ppt模板
    prs = Presentation(template)
    # 第一页
    slide_1 = prs.slides[0]
    shapes_1 = slide_1.shapes
    # 修改日期
    day = time.strftime("%Y-%m-%d", time.localtime())
    text_frame_12 = shapes_1[5].text_frame
    text_frame_12.clear()
    p_12 = text_frame_12.paragraphs[0]
    p_12.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    run_12 = p_12.add_run()
    run_12.text = day  # 内容
    font_12 = run_12.font
    font_12.name = 'Microsoft YaHei'  # 字体
    font_12.size = Pt(14)  # 字号
    font_12.bold = True  # 加粗
    # 第三页
    slide_3 = prs.slides[2]
    text_frame_31 = slide_3.shapes[4].text_frame
    text_frame_31.clear()
    p = text_frame_31.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_31 = p.add_run()
    run_31.text = word[5] + word[0]  # 内容
    font_31 = run_31.font
    font_31.name = 'Microsoft YaHei'  # 字体
    font_31.size = Pt(14)  # 字号
    # 第四页
    # 第一个文本框
    slide_4 = prs.slides[3]
    text_frame_41 = slide_4.shapes[4].text_frame
    text_frame_41.clear()
    p = text_frame_41.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_41 = p.add_run()
    run_41.text = word[1]  # 内容
    font_41 = run_41.font
    font_41.name = 'Microsoft YaHei'  # 字体
    font_41.size = Pt(14)  # 字号
    # 第二个文本框
    text_frame_42 = slide_4.shapes[5].text_frame
    text_frame_42.clear()
    p = text_frame_42.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_42 = p.add_run()
    run_42.text = word[2]  # 内容
    font_42 = run_42.font
    font_42.name = 'Microsoft YaHei'  # 字体
    font_42.size = Pt(14)  # 字号
    # 第五页-添加一张图片
    slide_5 = prs.slides[4]
    left, top, width, height = Cm(5.4), Cm(2.5), Cm(9), Cm(11)  # 预设位置及大小
    # 在指定位置按预设值添加图, width, height
    slide_5.shapes.add_picture(img_path, left, top)

    # 第六页-修改表格中的三个数值
    slide_6 = prs.slides[5]
    table = slide_6.shapes[4].table
    table.cell(1, 1).text = word[3]
    table.cell(1, 2).text = word[4]
    present = abs(float(word[3]) - float(word[4])) / float(word[3]) * 100
    table.cell(1, 3).text = str('%.2f' % present) + '%'
    try:
        prs.save(report_path)
        return 0
    except PermissionError:
        return report_path + '正在使用中，请关闭后重试！\n'


def create_ppt2(template, img_path, img_path2, report_path, word, word2):
    # 导入ppt模板
    prs = Presentation(template)
    # 第一页
    slide_1 = prs.slides[0]
    shapes_1 = slide_1.shapes
    # 修改日期
    day = time.strftime("%Y-%m-%d", time.localtime())
    text_frame_12 = shapes_1[5].text_frame
    text_frame_12.clear()
    p_12 = text_frame_12.paragraphs[0]
    p_12.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    run_12 = p_12.add_run()
    run_12.text = day  # 内容
    font_12 = run_12.font
    font_12.name = 'Microsoft YaHei'  # 字体
    font_12.size = Pt(14)  # 字号
    font_12.bold = True  # 加粗
    # 第三页
    # 第一个框
    slide_3 = prs.slides[2]
    text_frame_31 = slide_3.shapes[6].text_frame
    text_frame_31.clear()
    p = text_frame_31.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_31 = p.add_run()
    run_31.text = word[5] + word[0]  # 内容
    font_31 = run_31.font
    font_31.name = 'Microsoft YaHei'  # 字体
    font_31.size = Pt(14)  # 字号
    # font_31.bold = True
    # 第二个框
    text_frame_32 = slide_3.shapes[5].text_frame
    text_frame_32.clear()
    p = text_frame_32.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_32 = p.add_run()
    run_32.text = word2[5] + word2[0]  # 内容
    font_32 = run_32.font
    font_32.name = 'Microsoft YaHei'  # 字体
    font_32.size = Pt(14)  # 字号
    # font_32.bold = True
    # 第四页
    # 第一个文本框
    slide_4 = prs.slides[3]
    text_frame_41 = slide_4.shapes[4].text_frame
    text_frame_41.clear()
    p = text_frame_41.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_41 = p.add_run()
    run_41.text = word[1]  # 内容
    font_41 = run_41.font
    font_41.name = 'Microsoft YaHei'  # 字体
    font_41.size = Pt(14)  # 字号
    # 第二个文本框
    text_frame_42 = slide_4.shapes[5].text_frame
    text_frame_42.clear()
    p = text_frame_42.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_42 = p.add_run()
    run_42.text = word[2]  # 内容
    font_42 = run_42.font
    font_42.name = 'Microsoft YaHei'  # 字体
    font_42.size = Pt(14)  # 字号
    # 第三个文本框， 方案名
    text_frame_43 = slide_4.shapes[7].text_frame
    text_frame_43.clear()
    p = text_frame_43.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_43 = p.add_run()
    run_43.text = word[5]  # 内容
    font_43 = run_43.font
    font_43.name = 'Microsoft YaHei'  # 字体
    font_43.size = Pt(14)  # 字号

    # 第五页
    # 第一个文本框
    slide_5 = prs.slides[4]
    text_frame_51 = slide_5.shapes[4].text_frame
    text_frame_51.clear()
    p = text_frame_51.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_51 = p.add_run()
    run_51.text = word2[1]  # 内容
    font_51 = run_51.font
    font_51.name = 'Microsoft YaHei'  # 字体
    font_51.size = Pt(14)  # 字号
    # 第二个文本框
    text_frame_52 = slide_5.shapes[5].text_frame
    text_frame_52.clear()
    p = text_frame_52.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_52 = p.add_run()
    run_52.text = word2[2]  # 内容
    font_52 = run_52.font
    font_52.name = 'Microsoft YaHei'  # 字体
    font_52.size = Pt(14)  # 字号
    # 第三个文本框， 方案名
    text_frame_53 = slide_5.shapes[7].text_frame
    text_frame_53.clear()
    p = text_frame_53.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_53 = p.add_run()
    run_53.text = word2[5]  # 内容
    font_53 = run_53.font
    font_53.name = 'Microsoft YaHei'  # 字体
    font_53.size = Pt(14)  # 字号

    # 第六页-添加一张图片
    slide_6 = prs.slides[5]
    left, top, width, height = Cm(5.4), Cm(2.5), Cm(9), Cm(11)  # 预设位置及大小
    # 在指定位置按预设值添加图, width, height
    slide_6.shapes.add_picture(img_path, left, top)
    # 第七页-添加一张图片
    slide_7 = prs.slides[6]
    left, top, width, height = Cm(5.4), Cm(2.5), Cm(9), Cm(11)  # 预设位置及大小
    # 在指定位置按预设值添加图, width, height
    slide_7.shapes.add_picture(img_path2, left, top)

    # 第八页-修改表格中的三个数值
    slide_8 = prs.slides[7]
    table = slide_8.shapes[4].table
    # 方案1
    table.cell(1, 0).text = word[5].strip()
    table.cell(1, 1).text = word[3]
    table.cell(1, 2).text = word[4]
    present = abs(float(word[3]) - float(word[4])) / float(word[3]) * 100
    table.cell(1, 3).text = str('%.2f' % present) + '%'
    # 方案2
    table.cell(2, 0).text = word2[5].strip()
    table.cell(2, 1).text = word2[3]
    table.cell(2, 2).text = word2[4]
    present2 = abs(float(word2[3]) - float(word2[4])) / float(word2[3]) * 100
    table.cell(2, 3).text = str('%.2f' % present2) + '%'
    try:
        prs.save(report_path)
        return 0
    except PermissionError:
        return report_path + '正在使用中，请关闭后重试！\n'


def create_ppt3(
        template,
        img_path,
        img_path2,
        report_path,
        word,
        word2,
        word3):
    # 导入ppt模板
    prs = Presentation(template)
    # 第一页
    slide_1 = prs.slides[0]
    shapes_1 = slide_1.shapes
    # 修改日期
    day = time.strftime("%Y-%m-%d", time.localtime())
    text_frame_12 = shapes_1[5].text_frame
    text_frame_12.clear()
    p_12 = text_frame_12.paragraphs[0]
    p_12.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    run_12 = p_12.add_run()
    run_12.text = day  # 内容
    font_12 = run_12.font
    font_12.name = 'Microsoft YaHei'  # 字体
    font_12.size = Pt(14)  # 字号
    font_12.bold = True  # 加粗
    # 第三页
    # 第一个框
    slide_3 = prs.slides[2]
    text_frame_31 = slide_3.shapes[6].text_frame
    text_frame_31.clear()
    p = text_frame_31.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_31 = p.add_run()
    run_31.text = word[5] + word[0]  # 内容
    font_31 = run_31.font
    font_31.name = 'Microsoft YaHei'  # 字体
    font_31.size = Pt(14)  # 字号
    # font_31.bold = True
    # 第二个框
    text_frame_32 = slide_3.shapes[5].text_frame
    text_frame_32.clear()
    p = text_frame_32.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_32 = p.add_run()
    run_32.text = word2[5] + word2[0]  # 内容
    font_32 = run_32.font
    font_32.name = 'Microsoft YaHei'  # 字体
    font_32.size = Pt(14)  # 字号
    # font_32.bold = True
    # 第三个框
    text_frame_33 = slide_3.shapes[7].text_frame
    text_frame_33.clear()
    p = text_frame_33.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_33 = p.add_run()
    run_33.text = word3[5] + word3[0]  # 内容
    font_33 = run_33.font
    font_33.name = 'Microsoft YaHei'  # 字体
    font_33.size = Pt(14)  # 字号
    # font_33.bold = True
    # 第四页
    # 第一个文本框
    slide_4 = prs.slides[3]
    text_frame_41 = slide_4.shapes[4].text_frame
    text_frame_41.clear()
    p = text_frame_41.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_41 = p.add_run()
    run_41.text = word[1]  # 内容
    font_41 = run_41.font
    font_41.name = 'Microsoft YaHei'  # 字体
    font_41.size = Pt(14)  # 字号
    # 第二个文本框
    text_frame_42 = slide_4.shapes[5].text_frame
    text_frame_42.clear()
    p = text_frame_42.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_42 = p.add_run()
    run_42.text = word[2]  # 内容
    font_42 = run_42.font
    font_42.name = 'Microsoft YaHei'  # 字体
    font_42.size = Pt(14)  # 字号
    text_frame_43 = slide_4.shapes[7].text_frame
    text_frame_43.clear()
    p = text_frame_43.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_43 = p.add_run()
    run_43.text = word[5]  # 内容
    font_43 = run_43.font
    font_43.name = 'Microsoft YaHei'  # 字体
    font_43.size = Pt(14)  # 字号

    # 第五页
    # 第一个文本框
    slide_5 = prs.slides[4]
    text_frame_51 = slide_5.shapes[4].text_frame
    text_frame_51.clear()
    p = text_frame_51.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_51 = p.add_run()
    run_51.text = word2[1]  # 内容
    font_51 = run_51.font
    font_51.name = 'Microsoft YaHei'  # 字体
    font_51.size = Pt(14)  # 字号
    # 第二个文本框
    text_frame_52 = slide_5.shapes[5].text_frame
    text_frame_52.clear()
    p = text_frame_52.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_52 = p.add_run()
    run_52.text = word2[2]  # 内容
    font_52 = run_52.font
    font_52.name = 'Microsoft YaHei'  # 字体
    font_52.size = Pt(14)  # 字号
    # 第三个文本框， 方案名称
    text_frame_53 = slide_5.shapes[7].text_frame
    text_frame_53.clear()
    p = text_frame_53.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_53 = p.add_run()
    run_53.text = word2[5]  # 内容
    font_53 = run_53.font
    font_53.name = 'Microsoft YaHei'  # 字体
    font_53.size = Pt(14)  # 字号

    # 第六页
    # 第一个文本框
    slide_6 = prs.slides[5]
    text_frame_61 = slide_6.shapes[4].text_frame
    text_frame_61.clear()
    p = text_frame_61.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_61 = p.add_run()
    run_61.text = word3[1]  # 内容
    font_61 = run_61.font
    font_61.name = 'Microsoft YaHei'  # 字体
    font_61.size = Pt(14)  # 字号
    # 第二个文本框
    text_frame_62 = slide_6.shapes[5].text_frame
    text_frame_62.clear()
    p = text_frame_62.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_62 = p.add_run()
    run_62.text = word3[2]  # 内容
    font_62 = run_62.font
    font_62.name = 'Microsoft YaHei'  # 字体
    font_62.size = Pt(14)  # 字号
    # 第三个文本框， 方案名称
    text_frame_63 = slide_6.shapes[7].text_frame
    text_frame_63.clear()
    p = text_frame_63.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run_63 = p.add_run()
    run_63.text = word3[5]  # 内容
    font_63 = run_63.font
    font_63.name = 'Microsoft YaHei'  # 字体
    font_63.size = Pt(14)  # 字号

    # 第七页-添加一张图片
    slide_7 = prs.slides[6]
    left, top, width, height = Cm(5.4), Cm(2.5), Cm(9), Cm(11)  # 预设位置及大小
    # 在指定位置按预设值添加图, width, height
    slide_7.shapes.add_picture(img_path, left, top)
    # 第八页-添加一张图片
    slide_8 = prs.slides[7]
    left1, top1, width1, height1 = Cm(5.4), Cm(2.5), Cm(9), Cm(11)  # 预设位置及大小
    # 在指定位置按预设值添加图,
    slide_8.shapes.add_picture(img_path2, left1, top1)

    # 第九页-修改表格中的八个数值
    slide_9 = prs.slides[8]
    table = slide_9.shapes[4].table
    # 方案1
    table.cell(1, 0).text = word[5].strip()
    table.cell(1, 1).text = word[3]
    table.cell(1, 2).text = word[4]
    present = abs(float(word[3]) - float(word[4])) / float(word[3]) * 100
    table.cell(1, 3).text = str('%.2f' % present) + '%'
    # 方案2
    table.cell(2, 0).text = word2[5].strip()
    table.cell(2, 1).text = word2[3]
    table.cell(2, 2).text = word2[4]
    present2 = abs(float(word2[3]) - float(word2[4])) / float(word2[3]) * 100
    table.cell(2, 3).text = str('%.2f' % present2) + '%'
    # 方案3
    table.cell(3, 0).text = word3[5].strip()
    table.cell(3, 1).text = word3[3]
    table.cell(3, 2).text = word3[4]
    present3 = abs(float(word3[3]) - float(word3[4])) / float(word3[3]) * 100
    table.cell(3, 3).text = str('%.2f' % present3) + '%'
    try:
        prs.save(report_path)
        return 0
    except PermissionError:
        return report_path + '正在使用中，请关闭后重试！\n'


def main():
    message = ''
    # 读input文件
    if not os.path.exists('Input.dat'):
        return '未发现Input文件！\n'
    install_path, work_path, image_name, report_name, txt_name = read_input()
    if not install_path:
        return 'Input文件不完整，请检查！\n'

    template = install_path + '/_temp_.pptx'
    template2 = install_path + '/_temp_comparison.pptx'
    template3 = install_path + '/_temp_comparison2.pptx'
    if os.path.exists(template) and os.path.exists(template2) and os.path.exists(template3):
        if len(txt_name) == 2:
            # 创建方案1ppt
            txt_path = work_path + '/' + txt_name[1]
            word = read_txt(txt_path)
            img_path = work_path + '/' + image_name[1]
            report_path = work_path + '/' + report_name[1]
            msg1 = create_ppt(template, img_path, report_path, word)
            if isinstance(msg1, str):
                message += msg1
        # 如果存在方案2，则创建2个方案对比报告
        elif len(txt_name) == 3:
            img_path = work_path + '/' + image_name[1]
            img_path2 = work_path + '/' + image_name[2]
            report_path = work_path + '/' + report_name[1]
            txt_path = work_path + '/' + txt_name[1]
            txt_path2 = work_path + '/' + txt_name[2]
            word = read_txt(txt_path)
            word2 = read_txt(txt_path2)
            msg2 = create_ppt2(
                template2,
                img_path,
                img_path2,
                report_path,
                word,
                word2)
            if isinstance(msg2, str):
                message += msg2
        # 如果存在3个txt,导出3个的对比报告
        elif len(txt_name) == 4:
            img_path = work_path + '/' + image_name[1]
            img_path2 = work_path + '/' + image_name[2]
            report_path = work_path + '/' + report_name[1]
            txt_path = work_path + '/' + txt_name[1]
            txt_path2 = work_path + '/' + txt_name[2]
            txt_path3 = work_path + '/' + txt_name[3]
            word = read_txt(txt_path)
            word2 = read_txt(txt_path2)
            word3 = read_txt(txt_path3)
            # 导出3个报告的对比报告
            msg3 = create_ppt3(
                template3,
                img_path,
                img_path2,
                report_path,
                word,
                word2,
                word3)
            if isinstance(msg3, str):
                message += msg3
    else:
        return '模板文件缺失，请检查！\n'
    return message


if __name__ == '__main__':
    try:
        message1 = main()
    except Exception:
        with open('Output.dat', 'w') as out_f:
            out_f.write('导出失败!')
    else:
        with open('Output.dat', 'w') as out_f:
            if message1 == '':
                out_f.write('导出成功!')
            else:
                out_f.write(message1)
